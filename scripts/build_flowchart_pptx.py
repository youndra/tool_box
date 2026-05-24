#!/usr/bin/env python
"""Build an editable PPTX flowchart from a JSON specification."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_THEME = {
    "slide_width": 13.333,
    "slide_height": 7.5,
    "background_color": "#F7F4EA",
    "font_name": "Microsoft YaHei",
    "font_color": "#1F2937",
    "line_color": "#2F4858",
    "line_width": 1.5,
    "box_fill": "#FFFDF7",
    "box_line": "#2F4858",
    "accent_fill": "#E3B23C",
    "decision_fill": "#F4EBD0",
    "terminator_fill": "#D6EADF",
}

SHAPE_MAP = {
    "process": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "start": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
    "end": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
    "terminator": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
    "decision": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DECISION,
    "data": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DATA,
    "document": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT,
    "subprocess": MSO_AUTO_SHAPE_TYPE.FLOWCHART_PREDEFINED_PROCESS,
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
}

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create an editable PPTX flowchart from a JSON spec."
    )
    parser.add_argument("spec", help="Path to the JSON flowchart specification")
    parser.add_argument(
        "-o",
        "--output",
        help="Output PPTX path. Defaults to the input filename with .pptx",
    )
    return parser.parse_args()


def hex_to_rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    if len(value) != 6:
        raise ValueError(f"Expected 6-digit hex color, got: {value}")
    return RGBColor(*(int(value[i : i + 2], 16) for i in (0, 2, 4)))


def inches(value: float):
    return Inches(float(value))


def points(value: float):
    return Pt(float(value))


def merge_theme(spec_theme: dict | None) -> dict:
    theme = dict(DEFAULT_THEME)
    if spec_theme:
        theme.update(spec_theme)
    return theme


def get_shape_fill(node: dict, theme: dict) -> str:
    if "fill_color" in node:
        return node["fill_color"]
    shape_type = node.get("shape", "process")
    if shape_type in {"start", "end", "terminator"}:
        return theme["terminator_fill"]
    if shape_type == "decision":
        return theme["decision_fill"]
    return theme["box_fill"]


def apply_line_style(line, color_value, width_value, dash_value=None):
    if color_value == "none":
        line.fill.background()
    else:
        line.color.rgb = hex_to_rgb(color_value)
    if width_value is not None:
        line.width = points(width_value)
    if dash_value:
        if isinstance(dash_value, str):
            dash_map = {
                "dash": MSO_LINE_DASH_STYLE.DASH,
                "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
                "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
                "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
                "solid": MSO_LINE_DASH_STYLE.SOLID,
            }
            dash_value = dash_map[dash_value.lower()]
        line.dash_style = dash_value


def add_node(slide, node: dict, theme: dict):
    shape_kind = SHAPE_MAP.get(node.get("shape", "process"), SHAPE_MAP["process"])
    left = inches(node["x"])
    top = inches(node["y"])
    width = inches(node["w"])
    height = inches(node["h"])

    shape = slide.shapes.add_shape(shape_kind, left, top, width, height)
    fill_color = get_shape_fill(node, theme)
    if fill_color == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    apply_line_style(
        shape.line,
        node.get("line_color", theme["box_line"]),
        node.get("line_width", theme["line_width"]),
        node.get("line_dash"),
    )

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = node.get("text", "")
    font = run.font
    font.size = points(node.get("font_size", 16))
    font.bold = bool(node.get("bold", False))
    font.name = node.get("font_name", theme["font_name"])
    font.color.rgb = hex_to_rgb(node.get("font_color", theme["font_color"]))
    return shape


def anchor_point(node: dict, side: str) -> tuple[float, float]:
    x = float(node["x"])
    y = float(node["y"])
    w = float(node["w"])
    h = float(node["h"])
    side = side.lower()
    if side == "top":
        return x + w / 2, y
    if side == "right":
        return x + w, y + h / 2
    if side == "bottom":
        return x + w / 2, y + h
    if side == "left":
        return x, y + h / 2
    raise ValueError(f"Unsupported anchor side: {side}")


def add_label(slide, text: str, x: float, y: float, theme: dict):
    label = slide.shapes.add_textbox(inches(x), inches(y), inches(0.6), inches(0.25))
    text_frame = label.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.size = points(10)
    font.bold = True
    font.name = theme["font_name"]
    font.color.rgb = hex_to_rgb(theme["line_color"])
    label.line.fill.background()
    label.fill.background()


def add_free_text(slide, item: dict, theme: dict):
    box = slide.shapes.add_textbox(
        inches(item["x"]),
        inches(item["y"]),
        inches(item.get("w", 1.5)),
        inches(item.get("h", 0.4)),
    )
    text_frame = box.text_frame
    text_frame.clear()
    text_frame.word_wrap = bool(item.get("word_wrap", False))
    paragraph = text_frame.paragraphs[0]
    alignment = item.get("align", "center").lower()
    if alignment == "left":
        paragraph.alignment = PP_ALIGN.LEFT
    elif alignment == "right":
        paragraph.alignment = PP_ALIGN.RIGHT
    else:
        paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = item.get("text", "")
    font = run.font
    font.size = points(item.get("font_size", 16))
    font.bold = bool(item.get("bold", False))
    font.name = item.get("font_name", theme["font_name"])
    font.color.rgb = hex_to_rgb(item.get("font_color", theme["font_color"]))
    box.line.fill.background()
    box.fill.background()
    return box


def add_overlay_shape(slide, item: dict, theme: dict):
    shape_kind = SHAPE_MAP.get(item.get("shape", "rect"), SHAPE_MAP["rect"])
    shape = slide.shapes.add_shape(
        shape_kind,
        inches(item["x"]),
        inches(item["y"]),
        inches(item["w"]),
        inches(item["h"]),
    )
    fill_color = item.get("fill_color", "none")
    if fill_color == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    apply_line_style(
        shape.line,
        item.get("line_color", theme["line_color"]),
        item.get("line_width", theme["line_width"]),
        item.get("line_dash"),
    )
    if item.get("text"):
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = item["text"]
        font = run.font
        font.size = points(item.get("font_size", 14))
        font.bold = bool(item.get("bold", False))
        font.name = item.get("font_name", theme["font_name"])
        font.color.rgb = hex_to_rgb(item.get("font_color", theme["font_color"]))
    return shape


def add_segment(slide, item: dict, theme: dict):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        inches(item["x1"]),
        inches(item["y1"]),
        inches(item["x2"]),
        inches(item["y2"]),
    )
    apply_line_style(
        connector.line,
        item.get("line_color", theme["line_color"]),
        item.get("line_width", theme["line_width"]),
        item.get("dash"),
    )
    return connector


def add_connector(slide, edge: dict, nodes_by_id: dict, theme: dict):
    from_node = nodes_by_id[edge["from"]]
    to_node = nodes_by_id[edge["to"]]
    from_side = edge.get("from_side", "bottom")
    to_side = edge.get("to_side", "top")

    x1, y1 = anchor_point(from_node, from_side)
    x2, y2 = anchor_point(to_node, to_side)

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        inches(x1),
        inches(y1),
        inches(x2),
        inches(y2),
    )
    apply_line_style(
        connector.line,
        edge.get("line_color", theme["line_color"]),
        edge.get("line_width", theme["line_width"]),
        edge.get("dash"),
    )

    label = edge.get("label")
    if label:
        add_label(slide, label, (x1 + x2) / 2 - 0.3, (y1 + y2) / 2 - 0.125, theme)


def build_presentation(spec: dict, output_path: Path) -> None:
    theme = merge_theme(spec.get("theme"))
    prs = Presentation()
    prs.slide_width = inches(theme["slide_width"])
    prs.slide_height = inches(theme["slide_height"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(theme["background_color"])

    title = spec.get("title")
    if title:
        title_box = slide.shapes.add_textbox(inches(0.5), inches(0.2), inches(12.3), inches(0.5))
        title_frame = title_box.text_frame
        paragraph = title_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = points(24)
        run.font.name = theme["font_name"]
        run.font.color.rgb = hex_to_rgb(theme["font_color"])

    for overlay in spec.get("overlays", []):
        add_overlay_shape(slide, overlay, theme)

    for segment in spec.get("segments", []):
        add_segment(slide, segment, theme)

    nodes = spec.get("nodes", [])
    edges = spec.get("edges", [])

    nodes_by_id = {}
    for node in nodes:
        shape = add_node(slide, node, theme)
        node_copy = dict(node)
        node_copy["_shape"] = shape
        nodes_by_id[node["id"]] = node_copy

    for edge in edges:
        add_connector(slide, edge, nodes_by_id, theme)

    for text_item in spec.get("texts", []):
        add_free_text(slide, text_item, theme)

    prs.save(output_path)


def main() -> None:
    args = parse_args()
    spec_path = Path(args.spec)
    output_path = Path(args.output) if args.output else spec_path.with_suffix(".pptx")

    with spec_path.open("r", encoding="utf-8") as fh:
        spec = json.load(fh)

    build_presentation(spec, output_path)
    print(output_path)


if __name__ == "__main__":
    main()
